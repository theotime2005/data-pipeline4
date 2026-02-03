from pptx import Presentation
from pptx.util import Inches, Pt
import re

MD_PATH = 'slides/Presentation.md'
OUT_PPTX = 'slides/Presentation.pptx'

def split_slides(md_text):
    # split on lines that contain only ---
    parts = re.split(r"^---+$", md_text, flags=re.M)
    return [p.strip() for p in parts if p.strip()]

def clean_line(line):
    # remove leading markdown header markers
    return re.sub(r"^#+\s*", "", line).strip()

with open(MD_PATH, 'r', encoding='utf-8') as f:
    md = f.read()

slides_md = split_slides(md)
prs = Presentation()

first = True
for s in slides_md:
    lines = [ln.rstrip() for ln in s.splitlines() if ln.strip()]
    if not lines:
        continue
    title = clean_line(lines[0])
    body_lines = []
    for ln in lines[1:]:
        # list items starting with -, *, or numbered lists
        m = re.match(r"^[-*+]\s+(.*)$", ln)
        if m:
            body_lines.append(m.group(1))
            continue
        m2 = re.match(r"^\d+[.)]\s+(.*)$", ln)
        if m2:
            body_lines.append(m2.group(1))
            continue
        # code blocks or fenced blocks: skip fence markers
        if ln.startswith('```'):
            continue
        body_lines.append(clean_line(ln))

    # choose layout
    if first:
        # title slide
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        if slide.shapes.title:
            slide.shapes.title.text = title
        # subtitle: join first few body lines
        if len(body_lines) > 0:
            try:
                slide.placeholders[1].text = '\n'.join(body_lines[:4])
            except Exception:
                pass
        first = False
    else:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        if slide.shapes.title:
            slide.shapes.title.text = title
        # content placeholder
        try:
            body = slide.shapes.placeholders[1].text_frame
            for i, bl in enumerate(body_lines):
                if i == 0:
                    body.text = bl
                    p = body.paragraphs[0]
                    p.level = 0
                else:
                    p = body.add_paragraph()
                    p.text = bl
                    p.level = 0
        except Exception:
            # fallback: add a textbox
            left = Inches(1)
            top = Inches(1.5)
            width = Inches(8)
            height = Inches(4.5)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            for i, bl in enumerate(body_lines):
                if i == 0:
                    tf.text = bl
                else:
                    p = tf.add_paragraph()
                    p.text = bl

prs.save(OUT_PPTX)
print(f'Wrote {OUT_PPTX}')
